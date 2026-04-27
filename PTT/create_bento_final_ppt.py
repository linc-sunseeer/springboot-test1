from pathlib import Path
from shutil import copyfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt


BASE = Path(__file__).resolve().parent
OUT = BASE / "弁当予約管理システム_新人研修成果報告_最終版.pptx"
ROOT_OUT = BASE.parent / OUT.name

FONT = "Yu Gothic"
NAVY = RGBColor(28, 48, 70)
BLUE = RGBColor(48, 111, 178)
GREEN = RGBColor(47, 137, 98)
ORANGE = RGBColor(225, 139, 59)
RED = RGBColor(194, 76, 70)
BG = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(44, 52, 60)
MUTED = RGBColor(100, 111, 122)
LINE = RGBColor(212, 222, 228)


def style(run, size=16, bold=False, color=TEXT):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def textbox(slide, x, y, w, h, value, size=16, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.04)
    tf.margin_right = Cm(0.04)
    tf.margin_top = Cm(0.04)
    tf.margin_bottom = Cm(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = value
    style(r, size, bold, color)
    return box


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def header(slide, no, title, subtitle=""):
    textbox(slide, 0.9, 0.55, 1.2, 0.55, f"{no:02d}", 17, True, BLUE)
    textbox(slide, 2.0, 0.45, 22.8, 0.75, title, 27, True, NAVY)
    if subtitle:
        textbox(slide, 2.05, 1.25, 22.0, 0.35, subtitle, 11, False, MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0.9), Cm(1.85), Cm(24.8), Cm(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def new_slide(prs, no=None, title=None, subtitle=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    if title:
        header(s, no, title, subtitle)
    return s


def pill(slide, x, y, w, h, label, color=BLUE, size=12):
    sh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.adjustments[0] = 0.25
    tf = sh.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    style(r, size, True, WHITE)


def card(slide, x, y, w, h, title, items, color=BLUE, size=14):
    sh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = LINE
    sh.adjustments[0] = 0.06
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(x), Cm(y), Cm(0.13), Cm(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = sh.text_frame
    tf.clear()
    tf.margin_left = Cm(0.45)
    tf.margin_right = Cm(0.25)
    tf.margin_top = Cm(0.25)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    style(r, 16, True, NAVY)
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(4)


def table(slide, x, y, w, h, rows, widths=None, font_size=10):
    tb = slide.shapes.add_table(len(rows), len(rows[0]), Cm(x), Cm(y), Cm(w), Cm(h)).table
    if widths:
        for i, width in enumerate(widths):
            tb.columns[i].width = Cm(width)
    for ri, row in enumerate(rows):
        for ci, value in enumerate(row):
            cell = tb.cell(ri, ci)
            cell.text = str(value)
            cell.margin_left = Cm(0.08)
            cell.margin_right = Cm(0.08)
            cell.margin_top = Cm(0.05)
            cell.margin_bottom = Cm(0.05)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if ri == 0 else (RGBColor(240, 246, 250) if ri % 2 else WHITE)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if ri == 0 or ci == 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    style(r, font_size, ri == 0, WHITE if ri == 0 else TEXT)
    return tb


def arrow(slide, x1, y1, x2, y2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    c.line.color.rgb = MUTED
    c.line.width = Pt(2)
    c.line.end_arrowhead = True


prs = Presentation()
prs.slide_width = Cm(26.67)
prs.slide_height = Cm(15)

# 1
s = new_slide(prs)
left = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(6.2), Cm(15))
left.fill.solid()
left.fill.fore_color.rgb = NAVY
left.line.fill.background()
textbox(s, 0.8, 1.2, 4.8, 0.6, "株式会社サンシーア", 15, True, WHITE, PP_ALIGN.CENTER)
textbox(s, 0.8, 2.0, 4.8, 0.8, "新人研修", 22, True, WHITE, PP_ALIGN.CENTER)
textbox(s, 7.2, 4.1, 17.8, 1.0, "弁当予約管理システム", 34, True, NAVY)
textbox(s, 7.25, 5.35, 15.0, 0.65, "2026年度新人研修 成果報告", 22, True, BLUE)
textbox(s, 7.3, 6.55, 16.0, 0.8, "社内向けの予約・集計・管理をWebで行うシステム", 17, False, TEXT)
for i, (label, color) in enumerate([("Vue 3", GREEN), ("Spring Boot", BLUE), ("MyBatis-Plus", ORANGE), ("MySQL", RED)]):
    pill(s, 7.3 + i * 4.0, 8.25, 3.3, 0.62, label, color)

# 2
s = new_slide(prs, 1, "目次", "短く説明し、主要機能はWeb画面で実演")
agenda = ["概要", "開発環境", "ユーザー機能", "管理者機能", "デモの流れ", "データベース", "安全対策", "独自機能", "要件実現確認", "苦労した点", "改善点・課題", "まとめ"]
for i, item in enumerate(agenda):
    col = 0 if i < 6 else 1
    row = i if i < 6 else i - 6
    x = 2.1 + col * 11.7
    y = 2.7 + row * 1.35
    pill(s, x, y, 1.0, 0.58, str(i + 1), BLUE if col == 0 else GREEN, 10)
    textbox(s, x + 1.35, y + 0.08, 9.0, 0.35, item, 15, True)

# 3
s = new_slide(prs, 2, "システム概要", "弁当予約をWebで管理")
card(s, 1.3, 2.8, 11.5, 6.1, "目的", ["紙の予約管理をデジタル化", "予約状況をリアルタイムで確認", "集計作業を管理画面で効率化"], GREEN)
card(s, 14.0, 2.8, 11.5, 6.1, "利用者", ["一般社員：メニュー閲覧・予約", "管理者：メニュー登録・集計", "開発期間：2026年4月〜7月"], BLUE)
textbox(s, 2.0, 10.3, 22.8, 0.7, "発表では、スライド説明の後に実際の画面を操作して流れを確認します。", 17, True, NAVY, PP_ALIGN.CENTER)

# 4
s = new_slide(prs, 3, "開発環境", "実際のコード構成")
card(s, 1.2, 2.7, 11.5, 4.2, "バックエンド", ["Java 17", "Spring Boot 3.2", "Spring Security / JWT", "MyBatis-Plus"], BLUE, 13)
card(s, 14.0, 2.7, 11.5, 4.2, "フロントエンド", ["Vue 3", "Vite", "Pinia", "Vue Router / Tailwind CSS"], GREEN, 13)
card(s, 1.2, 8.0, 11.5, 3.2, "データベース", ["MySQL 8.0", "H2（テスト用）"], ORANGE, 13)
card(s, 14.0, 8.0, 11.5, 3.2, "実行環境", ["Docker Compose", "Frontend: 5174", "Backend: 8080"], RED, 13)

# 5
s = new_slide(prs, 4, "ユーザー機能", "一般社員が使う画面")
card(s, 1.2, 2.7, 7.5, 7.2, "アカウント", ["登録", "メール認証", "ログイン", "パスワード再設定"], BLUE, 13)
card(s, 9.55, 2.7, 7.5, 7.2, "予約", ["本日メニュー", "メニュー詳細", "カート予約", "予約履歴"], GREEN, 13)
card(s, 17.9, 2.7, 7.5, 7.2, "マイページ", ["氏名変更", "パスワード変更", "カロリー設定", "アレルゲン設定"], ORANGE, 13)
textbox(s, 2.0, 11.2, 22.8, 0.55, "この部分はWeb画面で、ログインから予約まで実演します。", 16, True, NAVY, PP_ALIGN.CENTER)

# 6
s = new_slide(prs, 5, "管理者機能", "管理者が使う画面")
card(s, 1.2, 2.7, 7.5, 7.2, "メニュー管理", ["登録・編集・削除", "画像アップロード", "提供日設定"], GREEN, 13)
card(s, 9.55, 2.7, 7.5, 7.2, "予約管理", ["予約状況確認", "支払確認", "成団確定", "返金処理"], BLUE, 13)
card(s, 17.9, 2.7, 7.5, 7.2, "集計", ["ダッシュボード", "日次集計", "月次集計", "CSV出力"], ORANGE, 13)
textbox(s, 2.0, 11.2, 22.8, 0.55, "管理画面では、メニュー登録と集計画面を中心に確認します。", 16, True, NAVY, PP_ALIGN.CENTER)

# 7
s = new_slide(prs, 6, "デモの流れ", "発表時に操作する順番")
flow = [("1", "ユーザー登録・ログイン", BLUE), ("2", "本日メニュー確認", GREEN), ("3", "カート予約", ORANGE), ("4", "マイページ確認", BLUE), ("5", "管理者ログイン", GREEN), ("6", "集計・メニュー管理", RED)]
for i, (num, label, color) in enumerate(flow):
    x = 1.2 + (i % 3) * 8.2
    y = 3.0 + (i // 3) * 3.3
    pill(s, x, y, 0.9, 0.65, num, color, 11)
    card(s, x + 1.1, y - 0.2, 6.2, 1.4, label, [], color, 12)
    if i in (0, 1, 3, 4):
        arrow(s, x + 7.35, y + 0.45, x + 8.0, y + 0.45)
textbox(s, 2.0, 10.9, 22.8, 0.6, "スライドは短くし、詳細は実際の画面で説明します。", 17, True, NAVY, PP_ALIGN.CENTER)

# 8
s = new_slide(prs, 7, "データベース", "主なテーブル")
table(s, 1.2, 2.6, 24.2, 6.7, [
    ["テーブル", "保存する内容"],
    ["users", "一般ユーザー、メール認証状態、食事設定"],
    ["admins", "管理者アカウント"],
    ["menus", "弁当メニュー、価格、画像、提供日、熱量、アレルゲン"],
    ["reservations", "予約、数量、支払状態、注文状態、返金理由"],
    ["email_tokens", "メール認証・パスワード再設定トークン"],
    ["system_settings", "予約締切時間、成団人数"],
], [5.0, 19.2], 11)
textbox(s, 2.0, 10.5, 22.8, 0.6, "予約は users と menus を reservations でつなぐ構成です。", 16, True, NAVY, PP_ALIGN.CENTER)

# 9
s = new_slide(prs, 8, "安全対策", "ログイン・権限・データ保護")
card(s, 1.2, 2.55, 11.5, 4.2, "ログイン管理", ["JWT Token を発行", "Bearer Token をAPIに送信", "FilterでTokenを検証"], BLUE, 13)
card(s, 14.0, 2.55, 11.5, 4.2, "ROLE判定", ["USER と ADMIN を分離", "一般画面：ROLE_USER", "管理画面：ROLE_ADMIN"], GREEN, 13)
card(s, 1.2, 7.8, 11.5, 3.5, "パスワード", ["BCryptでハッシュ化", "DBに平文保存しない"], ORANGE, 13)
card(s, 14.0, 7.8, 11.5, 3.5, "メール認証", ["登録後に本人確認", "未認証ユーザーのログインを防止"], RED, 13)

# 10
s = new_slide(prs, 9, "安全対策-API権限", "URLごとのアクセス制御")
table(s, 1.4, 3.0, 23.8, 5.2, [
    ["API", "権限", "内容"],
    ["/api/auth/**", "不要", "登録・ログイン・メール認証"],
    ["/api/user/**", "ROLE_USER / ROLE_ADMIN", "メニュー閲覧・予約・マイページ"],
    ["/api/admin/**", "ROLE_ADMIN", "メニュー管理・集計・設定"],
], [7.0, 7.2, 9.6], 12)
textbox(s, 2.0, 9.8, 22.8, 0.8, "管理者用APIは、ログインしていても ADMIN 権限がないと利用できません。", 17, True, NAVY, PP_ALIGN.CENTER)

# 11
s = new_slide(prs, 10, "独自機能", "要件に加えて工夫した機能")
card(s, 1.2, 2.7, 11.5, 4.0, "成団管理", ["最小予約人数を設定", "人数不足時のキャンセルに対応", "管理者が確定・取消を操作"], GREEN, 13)
card(s, 14.0, 2.7, 11.5, 4.0, "食事設定", ["目標カロリーを登録", "アレルゲンを登録", "メニュー情報と比較しやすい"], BLUE, 13)
card(s, 1.2, 7.8, 11.5, 3.5, "CSV出力", ["日次・月次集計を出力", "Excelで確認しやすい"], ORANGE, 13)
card(s, 14.0, 7.8, 11.5, 3.5, "返金処理", ["返金理由を記録", "予約履歴に状態を残す"], RED, 13)

# 12
s = new_slide(prs, 11, "要件実現確認", "要件定義に対する実装状況")
table(s, 1.0, 2.45, 24.6, 7.2, [
    ["要件", "状況", "確認方法"],
    ["ユーザー登録・ログイン", "実現済み", "登録画面・ログイン画面"],
    ["メニュー閲覧", "実現済み", "本日メニュー画面"],
    ["予約・履歴確認", "実現済み", "カート・マイページ"],
    ["メニュー管理", "実現済み", "管理者メニュー画面"],
    ["日次・月次集計", "実現済み", "管理者集計画面"],
    ["メール認証", "実現済み", "登録後の確認メール"],
], [7.0, 4.0, 13.6], 10)
textbox(s, 2.0, 10.7, 22.8, 0.6, "主要機能は実装済みで、発表では画面操作で確認します。", 16, True, NAVY, PP_ALIGN.CENTER)

# 13
s = new_slide(prs, 12, "苦労した点", "開発で難しかった部分")
card(s, 1.2, 2.8, 11.5, 6.4, "技術面", ["JWT と ROLE 判定", "メール送信設定", "画像アップロード", "フロントとAPIの連携"], RED, 13)
card(s, 14.0, 2.8, 11.5, 6.4, "開発面", ["Vue の学習", "画面数が多い", "テストデータの準備", "発表用に説明を整理"], ORANGE, 13)

# 14
s = new_slide(prs, 13, "改善点・課題", "今後さらに良くしたい点")
card(s, 1.2, 2.8, 11.5, 6.4, "機能", ["決済機能との連携", "スマホ画面の改善", "キャンセル操作の改善"], GREEN, 13)
card(s, 14.0, 2.8, 11.5, 6.4, "品質", ["画面操作を含むテスト追加", "デザインの統一", "パフォーマンス改善"], BLUE, 13)

# 15
s = new_slide(prs, 14, "まとめ", "今回の成果")
textbox(s, 2.0, 3.0, 22.8, 0.7, "Vue 3 + Spring Boot で弁当予約管理システムを作成", 23, True, NAVY, PP_ALIGN.CENTER)
card(s, 3.0, 4.7, 20.8, 4.6, "実装したこと", ["ユーザー予約機能", "管理者のメニュー・集計機能", "JWT / ROLE による安全対策", "成団管理・食事設定などの独自機能"], BLUE, 15)
textbox(s, 2.0, 10.6, 22.8, 0.6, "この後、実際のWeb画面で主要機能を確認します。", 18, True, GREEN, PP_ALIGN.CENTER)

# 16
s = new_slide(prs)
textbox(s, 2.0, 4.5, 22.6, 0.9, "ご清聴ありがとうございました", 36, True, NAVY, PP_ALIGN.CENTER)
textbox(s, 2.0, 6.2, 22.6, 0.7, "Q&A", 28, True, BLUE, PP_ALIGN.CENTER)

try:
    prs.save(OUT)
except PermissionError:
    OUT = BASE / "弁当予約管理システム_新人研修成果報告_最終版_発表用.pptx"
    ROOT_OUT = BASE.parent / OUT.name
    prs.save(OUT)

copyfile(OUT, ROOT_OUT)
print(OUT)
print(ROOT_OUT)
