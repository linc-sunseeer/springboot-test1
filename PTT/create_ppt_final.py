"""Generate FINAL PPT - simple overall + detailed unique features + speaking notes."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt

BASE_DIR = Path(__file__).resolve().parent
OUT = BASE_DIR / "弁当予約管理システム_新人研修成果報告_最終版.pptx"
ROOT_COPY = BASE_DIR.parent / "弁当予約管理システム_新人研修成果報告_最終版.pptx"

FONT = "Yu Gothic"
NAVY = RGBColor(24, 47, 72)
BLUE = RGBColor(55, 115, 185)
GREEN = RGBColor(45, 135, 95)
ORANGE = RGBColor(232, 149, 64)
RED = RGBColor(196, 78, 70)
TEXT = RGBColor(45, 52, 60)
MUTED = RGBColor(100, 111, 123)
BG = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(210, 220, 226)
NOTE_BG = RGBColor(255, 251, 235)
NOTE_BORDER = RGBColor(254, 215, 170)


def run_style(run, size=16, bold=False, color=TEXT):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def text(slide, x, y, w, h, value, size=16, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Cm(0.05)
    frame.margin_right = Cm(0.05)
    frame.margin_top = Cm(0.05)
    frame.margin_bottom = Cm(0.05)
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = value
    run_style(r, size, bold, color)
    return box


def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def make_slide(prs, title=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    if title:
        text(s, 1.0, 0.5, 24.0, 0.7, title, 26, True, NAVY)
        line = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(1.0), Cm(1.3), Cm(24.5), Cm(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = LINE
        line.line.fill.background()
    return s


def card(slide, x, y, w, h, title, items, color=BLUE, size=14):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = LINE
    s.adjustments[0] = 0.06
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(x), Cm(y), Cm(0.15), Cm(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    frame = s.text_frame
    frame.clear()
    frame.margin_left = Cm(0.5)
    frame.margin_right = Cm(0.3)
    frame.margin_top = Cm(0.35)
    p = frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    run_style(r, 16, True, NAVY)
    for item in items:
        p = frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(3)


def note(slide, x, y, w, h, text_value, size=11):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = NOTE_BG
    s.line.color.rgb = NOTE_BORDER
    s.adjustments[0] = 0.08
    frame = s.text_frame
    frame.clear()
    frame.margin_left = Cm(0.4)
    frame.margin_right = Cm(0.4)
    frame.margin_top = Cm(0.3)
    p = frame.paragraphs[0]
    r = p.add_run()
    r.text = text_value
    run_style(r, size, False, MUTED)


def pill(slide, x, y, w, h, value, color=BLUE, size=12):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.adjustments[0] = 0.25
    frame = s.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    run_style(r, size, True, WHITE)


prs = Presentation()
prs.slide_width = Cm(26.67)
prs.slide_height = Cm(15)

# 1 - Title
s = make_slide(prs)
left = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(6.2), Cm(15))
left.fill.solid()
left.fill.fore_color.rgb = NAVY
left.line.fill.background()
text(s, 0.8, 1.0, 4.8, 0.7, "株式会社サンシーア", 15, True, WHITE, PP_ALIGN.CENTER)
text(s, 0.8, 2.0, 4.8, 0.7, "新人研修", 22, True, WHITE, PP_ALIGN.CENTER)
text(s, 7.2, 4.0, 17.8, 1.0, "弁当予約管理システム", 34, True, NAVY)
text(s, 7.25, 5.25, 12.0, 0.7, "2026年度新人研修 成果報告", 22, True, BLUE)
text(s, 7.3, 6.55, 16.8, 1.2, "社内向け弁当予約・管理システム", 17, False, TEXT)
for i, (label, color) in enumerate([("Vue 3", GREEN), ("Spring Boot", BLUE), ("MyBatis-Plus", ORANGE)]):
    pill(s, 7.3 + i * 5.0, 8.5, 4.0, 0.7, label, color)

# 2 - Overview
s = make_slide(prs, "システム概要")
card(s, 1.2, 2.3, 11.5, 4.2, "開発目的", [
    "社内弁当予約をWeb化",
    "紙の管理からデジタル化",
    "予約状況をリアルタイム確認",
], GREEN)
card(s, 14.0, 2.3, 11.5, 4.2, "対象ユーザー", [
    "一般社員：予約・キャンセル",
    "管理者：メニュー管理・集計",
], BLUE)
text(s, 2.0, 7.2, 22.8, 0.5, "開発期間：2026年4月 〜 7月（約3ヶ月）", 15, True, MUTED, PP_ALIGN.CENTER)
note(s, 1.2, 8.2, 24.3, 1.5, "【口頭補足】前輩のシステムはThymeleafで画面とデータが一緒ですが、私はVue.jsを使って画面とAPIを分離しました。これで画面の動きが速くなり、スマホ対応もしやすくなります。", 11)

# 3 - Tech Stack
s = make_slide(prs, "開発環境")
card(s, 1.2, 2.3, 7.8, 4.8, "バックエンド", [
    "Java 17",
    "Spring Boot 3.2",
    "MyBatis-Plus",
    "Spring Security",
], BLUE, 14)
card(s, 9.5, 2.3, 7.8, 4.8, "フロントエンド", [
    "Vue 3",
    "Vite",
    "Pinia",
    "Tailwind CSS",
], GREEN, 14)
card(s, 17.8, 2.3, 7.8, 4.8, "その他", [
    "MySQL 8.0",
    "Docker",
    "Git / GitHub",
    "H2（テスト用）",
], ORANGE, 14)
note(s, 1.2, 7.8, 24.3, 1.8, "【口頭補足】MyBatis-PlusはSQLを簡単に書くことができるツールです。前輩はMyBatisを使っていましたが、私はPlus版を使ってコードを少なくしました。Spring Securityでログインの安全を管理しています。", 11)

# 4 - User Features
s = make_slide(prs, "ユーザー機能")
card(s, 1.2, 2.3, 11.5, 4.5, "アカウント", [
    "ユーザー登録・メール認証",
    "ログイン・パスワードリセット",
    "プロフィール変更",
], GREEN)
card(s, 14.0, 2.3, 11.5, 4.5, "予約", [
    "本日メニュー閲覧",
    "カートで複数予約",
    "予約履歴確認",
    "食事設定（カロリー・アレルゲン）",
], BLUE)
note(s, 1.2, 7.5, 24.3, 1.8, "【口頭補足】メール認証で本物のメールアドレスか確認します。パスワードを忘れたときはメールで再設定できます。カート機能で複数の弁当をまとめて予約できます。", 11)

# 5 - Admin Features
s = make_slide(prs, "管理者機能")
card(s, 1.2, 2.3, 11.5, 4.5, "メニュー管理", [
    "メニュー登録・編集・削除",
    "画像アップロード",
], GREEN)
card(s, 14.0, 2.3, 11.5, 4.5, "運用管理", [
    "予約状況確認",
    "日次・月次集計",
    "CSV出力",
    "ユーザー管理",
], BLUE)
note(s, 1.2, 7.5, 24.3, 1.8, "【口頭補足】管理者画面でメニューの写真をアップロードできます。日次・月次の予約数をグラフで確認できます。CSVボタンでExcelに出力して集計作業を楽にできます。", 11)

# 6 - Database
s = make_slide(prs, "データベース")
text(s, 1.2, 2.5, 24.0, 0.6, "主なテーブル（6つ）", 18, True, NAVY)
for i, (name, desc) in enumerate([
    ("users", "ユーザー情報"),
    ("admins", "管理者情報"),
    ("menus", "弁当メニュー"),
    ("reservations", "予約情報"),
    ("email_tokens", "認証トークン"),
    ("system_settings", "システム設定"),
]):
    col = i % 3
    row = i // 3
    x = 1.2 + col * 8.2
    y = 3.5 + row * 2.2
    pill(s, x, y, 3.5, 0.75, name, BLUE)
    text(s, x + 3.8, y + 0.15, 4.0, 0.5, desc, 13, False, TEXT)
note(s, 1.2, 8.2, 24.3, 1.5, "【口頭補足】usersテーブルに目標カロリーとアレルゲン情報を保存しています。reservationsテーブルで予約の状態（申請中・確定・キャンセル）を管理しています。", 11)

# 7 - Security (DETAILED)
s = make_slide(prs, "安全対策")
card(s, 1.2, 2.3, 7.8, 4.8, "認証", [
    "JWT Token でログイン管理",
    "メール認証で本人確認",
    "パスワードは暗号化保存",
], BLUE, 13)
card(s, 9.5, 2.3, 7.8, 4.8, "権限", [
    "一般ユーザーと管理者を分離",
    "管理者画面は特別なURL",
    "APIも権限チェックあり",
], GREEN, 13)
card(s, 17.8, 2.3, 7.8, 4.8, "その他", [
    "CSRF対策",
    "画像アップロードサイズ制限",
    "SQLインジェクション対策",
], ORANGE, 13)
note(s, 1.2, 7.8, 24.3, 2.2, "【口頭補足】JWTとは「ログインした状態を覚えておく仕組み」です。メール認証で「本当にこの人のメールアドレスか」を確認します。パスワードは暗号化（BCrypt）して保存するので、データベースを見ても元のパスワードはわかりません。", 11)

# 8 - Unique Features (DETAILED)
s = make_slide(prs, "独自機能")
card(s, 1.2, 2.3, 11.5, 4.5, "成団管理", [
    "最小予約人数を設定",
    "人数に達しないと注文不可",
    "管理者が手動で確定・キャンセル",
], GREEN)
card(s, 14.0, 2.3, 11.5, 4.5, "返金処理", [
    "管理者が返金ボタンで対応",
    "返金理由を記録",
    "履歴に残す",
], BLUE)
note(s, 1.2, 7.5, 24.3, 2.2, "【口頭補足】成団管理：例えば「10人以上予約しないと注文できない」という設定ができます。これで無駄な注文を防ぎます。返金処理：お金のやり取りは手動ですが、システムで「誰に・いくら・なぜ返金したか」を記録します。", 11)

# 9 - Diet Settings (DETAILED)
s = make_slide(prs, "食事設定")
card(s, 1.2, 2.3, 11.5, 4.5, "個人設定", [
    "目標カロリーを設定",
    "アレルゲン（アレルギー）を登録",
    "プロフィールで管理",
], GREEN)
card(s, 14.0, 2.3, 11.5, 4.5, "メニュー表示", [
    "メニューに熱量を表示",
    "アレルゲン情報を表示",
    "自分の設定と比較",
], BLUE)
note(s, 1.2, 7.5, 24.3, 2.2, "【口頭補足】例えば「エビアレルギーがある人」はプロフィールで「エビ」を登録します。メニューに「エビが入っています」と表示されるので、間違って注文しないようにできます。カロリーも表示して、ダイエット中の人が選びやすくなります。", 11)

# 10 - Challenges
s = make_slide(prs, "苦労した点")
card(s, 1.2, 2.3, 11.5, 4.5, "技術面", [
    "ログイン認証の設定",
    "メール送信の設定",
    "画像アップロード",
], RED)
card(s, 14.0, 2.3, 11.5, 4.5, "開発面", [
    "フロント・バック同時進行",
    "Vueの学習",
    "画面デザイン",
], ORANGE)
note(s, 1.2, 7.5, 24.3, 1.8, "【口頭補足】JWTの設定が一番難しかったです。ネットで調べながら1週間かかりました。メール送信もGmailの設定でエラーが出て、先輩に聞いて解決しました。", 11)

# 11 - Improvements
s = make_slide(prs, "改善点・課題")
card(s, 1.2, 2.3, 11.5, 4.5, "機能", [
    "決済機能の追加",
    "スマホ対応",
    "キャンセル機能の改善",
], GREEN)
card(s, 14.0, 2.3, 11.5, 4.5, "技術", [
    "デザイン改善",
    "テスト追加",
    "パフォーマンス改善",
], BLUE)
note(s, 1.2, 7.5, 24.3, 1.8, "【口頭補足】今はお金のやり取りは手動ですが、PayPayやクレジットカードと連携したいです。スマホで使いやすい画面に改善したいです。", 11)

# 12 - Summary
s = make_slide(prs, "まとめ")
text(s, 2.0, 2.8, 22.6, 0.8, "3ヶ月で弁当予約システムを完成", 24, True, NAVY, PP_ALIGN.CENTER)
text(s, 2.0, 4.0, 22.6, 2.0, "・Vue 3 + Spring Boot で前後端分離\n・メール認証・画像アップロードなど実装\n・テスト 18 件実施\n・成団管理・食事設定など独自機能あり", 16, False, TEXT, PP_ALIGN.CENTER)
note(s, 1.2, 7.5, 24.3, 1.8, "【口頭補足】前輩のシステムは画面とデータが一緒でしたが、私は最新の技術で画面とAPIを分離しました。メール認証や画像アップロードなど、前輩にはない機能も作りました。", 11)

# 13 - Q&A
s = make_slide(prs)
bg(s)
text(s, 2.0, 4.8, 22.6, 1.0, "ご清聴ありがとうございました", 34, True, NAVY, PP_ALIGN.CENTER)
text(s, 2.0, 6.2, 22.6, 0.7, "Q&A", 28, True, BLUE, PP_ALIGN.CENTER)

prs.save(OUT)
ROOT_COPY.write_bytes(OUT.read_bytes())
print(f"Generated: {OUT}")
print(f"Copied to: {ROOT_COPY}")
