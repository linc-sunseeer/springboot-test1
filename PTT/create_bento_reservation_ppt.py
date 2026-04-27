from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt


BASE_DIR = Path(__file__).resolve().parent
OUT = BASE_DIR / "弁当予約管理システム_新人研修成果報告.pptx"
ROOT_COPY = BASE_DIR.parent / "弁当予約管理システム_新人研修成果報告.pptx"

FONT = "Yu Gothic"
NAVY = RGBColor(24, 47, 72)
BLUE = RGBColor(55, 115, 185)
GREEN = RGBColor(45, 135, 95)
ORANGE = RGBColor(232, 149, 64)
RED = RGBColor(196, 78, 70)
TEXT = RGBColor(45, 52, 60)
MUTED = RGBColor(100, 111, 123)
BG = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(210, 220, 226)


def run_style(run, size=16, bold=False, color=TEXT):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def text(slide, x, y, w, h, value, size=16, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Cm(0.05)
    frame.margin_right = Cm(0.05)
    frame.margin_top = Cm(0.05)
    frame.margin_bottom = Cm(0.05)
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = value
    run_style(r, size, bold, color)
    return box


def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def header(slide, no, title, subtitle=None):
    text(slide, 0.9, 0.55, 1.1, 0.65, f"{no:02d}", 18, True, BLUE)
    text(slide, 2.0, 0.5, 22.5, 0.8, title, 27, True, NAVY)
    if subtitle:
        text(slide, 2.05, 1.35, 20.5, 0.4, subtitle, 11, False, MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0.9), Cm(1.95), Cm(24.8), Cm(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def slide(prs, no=None, title=None, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    if title:
        header(s, no, title, subtitle)
    return s


def card(slide, x, y, w, h, title, items, color=BLUE, size=14):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = LINE
    s.adjustments[0] = 0.06
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(x), Cm(y), Cm(0.13), Cm(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    frame = s.text_frame
    frame.clear()
    frame.margin_left = Cm(0.45)
    frame.margin_right = Cm(0.25)
    frame.margin_top = Cm(0.3)
    p = frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    run_style(r, 16, True, NAVY)
    for item in items:
        p = frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(3)


def pill(slide, x, y, w, h, value, color=BLUE, size=12):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.adjustments[0] = 0.25
    frame = s.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    run_style(r, size, True, WHITE)


def add_table(slide, x, y, w, h, rows, widths=None, font_size=10):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Cm(x), Cm(y), Cm(w), Cm(h)).table
    if widths:
        for i, width in enumerate(widths):
            table.columns[i].width = Cm(width)
    for ri, row in enumerate(rows):
        for ci, value in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(value)
            cell.margin_left = Cm(0.08)
            cell.margin_right = Cm(0.08)
            cell.margin_top = Cm(0.05)
            cell.margin_bottom = Cm(0.05)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if ri == 0 else (RGBColor(240, 246, 250) if ri % 2 else WHITE)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if ri == 0 or ci == 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    run_style(r, font_size, ri == 0, WHITE if ri == 0 else TEXT)
    return table


def arrow(slide, x1, y1, x2, y2, color=MUTED):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    c.line.color.rgb = color
    c.line.width = Pt(2)
    c.line.end_arrowhead = True


prs = Presentation()
prs.slide_width = Cm(26.67)
prs.slide_height = Cm(15)

# 1
s = slide(prs)
left = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(6.2), Cm(15))
left.fill.solid()
left.fill.fore_color.rgb = NAVY
left.line.fill.background()
text(s, 0.8, 1.0, 4.8, 0.7, "株式会社サンシーア", 15, True, WHITE, PP_ALIGN.CENTER)
text(s, 0.8, 2.0, 4.8, 0.7, "新人研修", 22, True, WHITE, PP_ALIGN.CENTER)
text(s, 7.2, 4.0, 17.8, 1.0, "弁当予約管理システム", 34, True, NAVY)
text(s, 7.25, 5.25, 12.0, 0.7, "2026年度新人研修 成果報告", 22, True, BLUE)
text(s, 7.3, 6.55, 16.8, 1.2, "Vue 3 と Spring Boot による、社内向け弁当予約・集計管理システム", 17, False, TEXT)
for i, (label, color) in enumerate([("Vue 3", GREEN), ("Spring Boot", BLUE), ("MyBatis-Plus", ORANGE), ("MySQL 8.0", RED)]):
    pill(s, 7.3 + i * 4.1, 8.5, 3.4, 0.65, label, color)

# 2
s = slide(prs, 1, "システム概要", "実装済みコードに基づく概要")
card(s, 1.3, 2.75, 11.5, 7.5, "目的", [
    "社内の弁当予約をWeb化",
    "予約締切・成団条件をシステムで管理",
    "日次・月次集計を管理者画面で確認",
    "CSV出力により集計作業を効率化",
], GREEN)
card(s, 14.0, 2.75, 11.5, 7.5, "対象ユーザー", [
    "一般ユーザー：登録、メール認証、ログイン",
    "一般ユーザー：メニュー閲覧、カート予約、支払申請",
    "管理者：メニュー管理、予約集計、ユーザー管理",
    "管理者：締切時間・成団人数の設定",
], BLUE)
text(s, 2.0, 11.3, 22.8, 0.8, "紙運用ではなく、予約・支払・集計・通知をAPI中心に扱う構成", 17, True, NAVY, PP_ALIGN.CENTER)

# 3
s = slide(prs, 2, "機能一覧", "Controller / Router から確認できる実装範囲")
card(s, 1.1, 2.6, 7.7, 8.7, "認証・ユーザー", [
    "ユーザー登録",
    "メール認証",
    "JWTログイン",
    "パスワード再設定",
    "プロフィール変更",
    "目標カロリー・アレルゲン設定",
], GREEN, 12)
card(s, 9.45, 2.6, 7.7, 8.7, "予約機能", [
    "本日のメニュー取得",
    "週間メニュー表示",
    "メニュー詳細表示",
    "カート一括予約",
    "予約履歴確認",
    "支払申請・注文状態更新",
], BLUE, 12)
card(s, 17.8, 2.6, 7.7, 8.7, "管理者機能", [
    "ダッシュボード統計",
    "メニューCRUD・画像アップロード",
    "日次・月次集計",
    "CSVエクスポート",
    "予約確定・全件キャンセル",
    "返金・ユーザー管理",
], ORANGE, 12)

# 4
s = slide(prs, 3, "開発環境＆開発言語", "pom.xml / package.json / compose.yaml に基づく")
card(s, 1.2, 2.55, 11.6, 4.3, "バックエンド", [
    "Java 17",
    "Spring Boot 3.2.6 / Maven",
    "Spring Security / JWT",
    "MyBatis-Plus / Lombok",
    "Spring Mail",
], BLUE, 12)
card(s, 14.0, 2.55, 11.6, 4.3, "フロントエンド", [
    "Vue 3",
    "Vite",
    "Pinia",
    "Vue Router",
    "Tailwind CSS",
], GREEN, 12)
card(s, 1.2, 7.6, 11.6, 3.7, "データベース・テスト", [
    "MySQL 8.0",
    "H2（テスト用）",
    "schema.sql / data.sql",
], ORANGE, 12)
card(s, 14.0, 7.6, 11.6, 3.7, "リリース構成", [
    "Docker / Docker Compose",
    "backend: 8080",
    "frontend: 5174",
    "mysql: 3308 -> 3306",
], RED, 12)

# 5
s = slide(prs, 4, "開発スケジュール", "前輩資料の形式に合わせた工程整理")
schedule = [
    ("4月", "要件整理・DB設計", GREEN),
    ("5月", "API設計・認証実装", BLUE),
    ("6月", "予約・管理画面実装", ORANGE),
    ("7月", "テスト・発表準備", RED),
]
for i, (month, task, color) in enumerate(schedule):
    x = 1.6 + i * 6.0
    pill(s, x, 3.2, 4.5, 0.75, month, color, 15)
    card(s, x, 4.35, 4.5, 3.4, task, ["設計", "実装", "確認"][0:1], color, 12)
    text(s, x + 0.25, 5.55, 4.0, 0.8, task, 14, True, TEXT, PP_ALIGN.CENTER)
    if i < 3:
        arrow(s, x + 4.65, 5.9, x + 5.75, 5.9)
add_table(s, 2.0, 9.2, 22.4, 2.8, [["期間", "作業内容"], ["4月", "要件整理・DB設計"], ["5月", "API設計・認証実装"], ["6月", "予約・管理画面実装"], ["7月", "テスト・発表準備"]], [4, 18.4], 12)

# 6
s = slide(prs, 5, "基本設計-アーキテクチャ図", "Vue SPA + REST API + MyBatis-Plus")
for x, label, color in [(1.5, "Vue 3\n画面", GREEN), (7.0, "REST API\nController", BLUE), (12.7, "Service\n業務処理", ORANGE), (18.1, "Mapper\nMyBatis-Plus", BLUE), (22.3, "MySQL\nDB", RED)]:
    pill(s, x, 5.2, 3.5, 1.3, label, color, 12)
for x1, x2 in [(5.1, 6.9), (10.6, 12.6), (16.3, 18.0), (21.7, 22.2)]:
    arrow(s, x1, 5.85, x2, 5.85)
card(s, 2.0, 8.0, 10.5, 3.0, "認証", ["Spring Security はステートレス設定", "Bearer Token Filter で JWT を検証", "/api/admin/** は ADMIN 権限が必要"], BLUE, 12)
card(s, 14.0, 8.0, 10.5, 3.0, "メール", ["登録時のメール認証", "パスワード再設定メール", "予約リマインダー / 成団通知"], GREEN, 12)

# 7
s = slide(prs, 6, "基本設計-データベース", "schema.sql の主要テーブル")
add_table(s, 1.1, 2.65, 24.4, 5.0, [
    ["テーブル名", "役割"],
    ["users", "一般ユーザー情報、メール認証状態、カロリー・アレルゲン設定"],
    ["admins", "管理者ログイン情報"],
    ["menus", "弁当メニュー、価格、画像、提供日、熱量、アレルゲン"],
    ["reservations", "予約情報、数量、支払状態、注文状態、返金理由"],
    ["email_tokens", "メール認証・パスワード再設定トークン"],
    ["system_settings", "予約締切時間、成団人数などの設定"],
], [5, 19.4], 11)
for x, y, name, color in [
    (2.0, 9.25, "users", GREEN),
    (7.0, 9.25, "reservations", BLUE),
    (13.2, 9.25, "menus", ORANGE),
    (18.7, 9.25, "system_settings", RED),
]:
    pill(s, x, y, 4.2, 0.8, name, color)
arrow(s, 6.25, 9.65, 6.95, 9.65)
arrow(s, 11.3, 9.65, 13.1, 9.65)

# 8
s = slide(prs, 7, "基本設計-画面遷移図", "frontend/src/router/index.js に基づく")
user_flow = ["ログイン", "本日のメニュー", "メニュー詳細", "カート", "マイページ"]
for i, label in enumerate(user_flow):
    x = 1.2 + i * 4.8
    pill(s, x, 3.4, 3.7, 0.82, label, GREEN if i else BLUE, 11)
    if i < len(user_flow) - 1:
        arrow(s, x + 3.75, 3.8, x + 4.65, 3.8)
text(s, 1.25, 2.75, 7.0, 0.4, "一般ユーザー", 14, True, NAVY)
admin_flow = ["管理者ログイン", "ダッシュボード", "メニュー管理", "日次集計", "月次集計", "設定・ユーザー"]
for i, label in enumerate(admin_flow):
    x = 1.2 + (i % 3) * 7.5
    y = 7.2 + (i // 3) * 2.0
    pill(s, x, y, 5.4, 0.82, label, BLUE if i < 2 else ORANGE, 11)
    if i in (0, 1):
        arrow(s, x + 5.45, y + 0.4, x + 7.35, y + 0.4)
text(s, 1.25, 6.55, 7.0, 0.4, "管理者", 14, True, NAVY)

# 9
s = slide(prs, 8, "画面説明", "主要画面とできること")
card(s, 1.2, 2.6, 11.5, 3.5, "本日のメニュー画面", ["予約可能メニュー、締切時間、成団人数を表示", "メニュー詳細・カート画面へ遷移"], GREEN, 12)
card(s, 14.0, 2.6, 11.5, 3.5, "カート・予約画面", ["複数メニューをまとめて予約", "支払方法・数量を指定"], BLUE, 12)
card(s, 1.2, 7.2, 11.5, 3.5, "マイページ", ["予約履歴、支払申請、注文状態確認", "氏名・パスワード・食事設定を変更"], ORANGE, 12)
card(s, 14.0, 7.2, 11.5, 3.5, "管理画面", ["集計、メニュー管理、画像アップロード", "CSV出力、返金、成団確定を操作"], RED, 12)

# 10
s = slide(prs, 9, "テスト仕様書", "src/test/java にあるテストを整理")
add_table(s, 1.0, 2.55, 24.8, 7.0, [
    ["分類", "確認内容", "対象テスト"],
    ["認証", "ログイン、メール認証、パスワード再設定", "AuthControllerMailFlowTest / PasswordResetServiceTest"],
    ["セキュリティ", "APIの権限、JWT認証、管理者制御", "ApiSkeletonSecurityTest"],
    ["予約", "予約作成、カート予約、支払状態、注文状態", "ReservationServiceTest / ReservationControllerTest"],
    ["メニュー", "本日メニュー取得、メニュー管理", "MenuServiceTest / UserMenuControllerTest"],
    ["管理", "日次・月次集計、ユーザー管理", "AdminAggregateServiceTest / AdminControllerTest"],
    ["メール", "送信ON/OFF、通知処理", "MailServiceTest / EmailVerificationServiceTest"],
], [4.2, 10.3, 10.3], 10)
text(s, 2.0, 10.7, 22.8, 0.6, "単体テストだけでなく、Controller 層のAPI動作も確認", 16, True, NAVY, PP_ALIGN.CENTER)

# 11
s = slide(prs, 10, "苦労した点", "実装内容から見た技術課題")
card(s, 1.4, 2.75, 11.6, 7.4, "バックエンド", [
    "Spring Security と JWT の権限制御",
    "メール認証・パスワード再設定のトークン管理",
    "予約締切・成団条件・状態更新の整理",
    "CSV出力と集計ロジック",
], RED, 13)
card(s, 14.0, 2.75, 11.6, 7.4, "フロントエンド・連携", [
    "Vue Router の権限別画面遷移",
    "Pinia でのログイン状態管理",
    "管理者画面とユーザー画面のAPI分離",
    "画像アップロードと表示の扱い",
], ORANGE, 13)

# 12
s = slide(prs, 11, "改善点・課題", "今後の品質向上")
card(s, 1.4, 2.75, 11.6, 7.4, "機能面", [
    "予約キャンセル操作の画面上での分かりやすさ向上",
    "支払方法の外部決済連携",
    "メニュー画像・アレルゲン表示の改善",
    "スマートフォン画面の操作性向上",
], GREEN, 13)
card(s, 14.0, 2.75, 11.6, 7.4, "技術面", [
    "トランザクション境界の整理",
    "APIレスポンス形式の統一徹底",
    "E2Eテストの追加",
    "本番用メール・JWT secret 管理の強化",
], BLUE, 13)

# 13
s = slide(prs)
bg(s)
text(s, 2.0, 4.5, 22.6, 0.9, "ご清聴ありがとうございました", 36, True, NAVY, PP_ALIGN.CENTER)
text(s, 2.0, 6.2, 22.6, 0.7, "Q&A", 28, True, BLUE, PP_ALIGN.CENTER)
text(s, 2.0, 8.0, 22.6, 0.5, "弁当予約管理システム / 2026年度新人研修 成果報告", 14, False, MUTED, PP_ALIGN.CENTER)

prs.save(OUT)
ROOT_COPY.write_bytes(OUT.read_bytes())
print(OUT)
print(ROOT_COPY)
